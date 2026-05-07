from __future__ import annotations

from pathlib import Path

from scraper.office_text_extraction import (
    extract_office_text_from_bytes,
    is_binary_like_text,
    is_office_document,
)


def test_extract_office_text_from_bytes_uses_ipfs_powerpoint_extractor(tmp_path: Path):
    from pptx import Presentation

    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "211info child care referral toolkit"
    textbox = slide.shapes.add_textbox(100, 100, 400, 100)
    textbox.text = "Call 211 for child care navigation"
    presentation.save(pptx_path)

    result = extract_office_text_from_bytes(pptx_path.read_bytes(), source_name="sample.pptx")

    assert result.success
    assert "211info child care referral toolkit" in result.text
    assert "child care navigation" in result.text
    assert result.byte_length == pptx_path.stat().st_size
    assert result.binary_sha256


def test_is_office_document_and_binary_text_detection():
    assert is_office_document("https://example.org/toolkit.pptx")
    assert is_office_document(
        "https://example.org/file",
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    assert is_office_document(content=b"PK\x03\x04")
    assert is_binary_like_text("PK\x03\x04\x14\x00\x00\x00[Content_Types].xml")
    assert not is_binary_like_text("Plain service navigation text for a food pantry.")
